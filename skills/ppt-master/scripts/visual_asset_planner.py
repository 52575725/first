#!/usr/bin/env python3
"""Plan, extract, search, and generate visual assets for a PPT Master project.

The planner answers four questions per slide/paragraph:
1. Does this content need an image?
2. Can an existing image from the source PPT be reused?
3. If not, should the image be searched from licensed web sources or generated?
4. Is the selected image actually usable (readable, large enough, sane aspect)?

Default mode is safe: it extracts and validates existing assets, then writes a
dry-run plan. Use ``--execute`` to run web search / image generation commands.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import math
import os
import re
import shutil
import subprocess
import sys
import zipfile
from dataclasses import dataclass, asdict, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable


SCRIPT_DIR = Path(__file__).resolve().parent
REPO_ROOT = SCRIPT_DIR.parent.parent.parent
VENDOR_DIR = REPO_ROOT / "tools" / "python_libs"
if VENDOR_DIR.exists() and str(VENDOR_DIR) not in sys.path:
    sys.path.insert(0, str(VENDOR_DIR))
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif", ".tif", ".tiff"}
PPT_EXTENSIONS = {".pptx", ".pptm", ".ppsx", ".ppsm", ".potx", ".potm"}

try:
    from PIL import Image, ImageStat  # type: ignore
except ImportError:  # pragma: no cover - runtime fallback
    Image = None
    ImageStat = None

try:
    from subject_framework import framework_for_slide
except Exception:  # pragma: no cover - keep dry-run planner usable standalone
    framework_for_slide = None


CONCRETE_KEYWORDS = (
    "旧衣", "衣物", "纺织", "面料", "服装", "垃圾", "填埋", "焚烧", "回收箱",
    "分拣", "清洗", "消毒", "纤维", "工厂", "设备", "门店", "城市", "社区",
    "团队", "消费者", "品牌", "产品", "材料", "案例", "场景", "步骤", "流程",
    "结构", "装置", "实验", "地图", "人物", "照片", "示意图", "图像", "函数",
    "图形", "根号", "平方根", "二次根式", "数轴", "三角形", "面积", "边长",
    "速度", "建筑", "工程", "用户", "竞赛", "创新", "创业", "fabric", "textile",
    "recycling", "factory", "team", "product", "store", "city", "landfill",
    "diagram", "example", "case", "process", "workflow", "map", "function",
    "graph", "geometry", "语文", "文学", "课文", "文本", "意象", "修辞", "作者",
    "史铁生", "我与地坛", "地坛", "历史", "朝代", "革命", "战争", "地理", "气候",
    "经纬度", "等高线", "摩擦力", "受力分析", "测力计", "电路", "化学", "细胞",
    "DNA", "算法", "代码",
)

ABSTRACT_KEYWORDS = (
    "战略", "模式", "价值", "愿景", "使命", "算法", "平台", "体系", "模型",
    "利润", "收入", "市场规模", "政策", "优势", "壁垒", "规划", "定义", "概念",
    "原理", "规律", "方法", "关系", "roadmap", "strategy", "model", "policy",
    "revenue", "profit", "advantage",
)

NO_IMAGE_KEYWORDS = (
    "目录", "agenda", "数据", "指标", "测算", "KPI", "ROI", "财务", "表格",
)

SEARCH_QUERY_RULES = (
    (("根号", "定义", "平方根"), "square root definition diagram"),
    (("算术平方根", "平方根"), "arithmetic square root diagram"),
    (("根号", "性质", "非负"), "square root nonnegative domain diagram"),
    (("根号", "分类", "立方根"), "square root cube root comparison diagram"),
    (("乘法", "根号"), "square root multiplication rule diagram"),
    (("除法", "根号"), "square root division rule diagram"),
    (("加减", "同类根式"), "like radicals addition subtraction diagram"),
    (("化简", "完全平方因子"), "simplifying radicals perfect square factors diagram"),
    (("无理数", "根号"), "irrational numbers square root number line diagram"),
    (("方程", "根号"), "solving radical equations example diagram"),
    (("不等式", "根号", "定义域"), "radical inequality domain diagram"),
    (("金融", "复利"), "compound interest formula growth diagram"),
    (("物理", "速度"), "physics velocity formula square root diagram"),
    (("万有引力", "距离"), "physics gravity distance formula diagram"),
    (("函数", "单调性"), "square root function monotonicity graph"),
    (("函数", "图像"), "square root function graph"),
    (("根号", "平方根", "开方"), "square root concept diagram"),
    (("二次根式", "根式"), "radical expression math diagram"),
    (("勾股", "直角三角形", "三角形"), "Pythagorean theorem diagram"),
    (("定义域", "数轴"), "number line square root domain diagram"),
    (("函数", "曲线", "图像"), "square root function graph"),
    (("例题", "解题", "计算"), "math problem solving example"),
    (("面积", "边长", "正方形"), "square area side length square root diagram"),
    (("速度", "加速度"), "physics velocity acceleration diagram"),
    (("填埋", "垃圾", "废料"), "textile waste landfill"),
    (("焚烧", "污染"), "clothing waste pollution"),
    (("智能回收", "回收箱"), "clothing recycling bin"),
    (("AI", "视觉", "分拣"), "textile sorting machine"),
    (("分拣",), "textile sorting facility"),
    (("清洗", "消毒"), "textile washing factory"),
    (("再生纤维", "纤维"), "recycled textile fiber"),
    (("环保面料", "面料"), "recycled fabric textile"),
    (("市场", "消费"), "sustainable fashion store"),
    (("品牌",), "sustainable fashion retail"),
    (("团队", "CEO", "CTO", "COO"), "business team meeting"),
    (("全国", "城市", "布局"), "China city skyline"),
    (("旧衣", "衣物", "纺织"), "textile recycling"),
)

RELEVANCE_ACCEPT_THRESHOLD = 0.34
WEB_RELEVANCE_ACCEPT_THRESHOLD = 0.30
VISUAL_CACHE_DIR = REPO_ROOT / "assets" / "visual_cache"
GENERIC_ASSET_TOKENS = {
    "asset", "assets", "candidate", "cover", "diagram", "extracted", "file", "files",
    "image", "img", "input", "media", "photo", "pic", "picture", "ppt", "pptx", "shape",
    "slide", "source", "src", "tmp", "visual", "webp", "jpeg", "jpg", "png",
    "editorial", "illustration", "landscape", "portrait", "background", "high", "quality",
    "clean", "education", "educational",
}
ENGLISH_STOPWORDS = {
    "about", "after", "and", "are", "can", "for", "from", "how", "into", "its", "more",
    "not", "of", "on", "our", "page", "part", "ppt", "slide", "the", "this", "that",
    "their", "they", "use", "using", "with", "you",
}
CHINESE_STOPWORDS = {
    "这一页", "本页", "页面", "内容", "介绍", "说明", "什么是", "我们", "可以", "进行",
    "通过", "理解", "掌握", "知道", "分析", "展示", "问题", "方法",
}
WEB_SUBJECT_MISMATCH_TERMS = {
    "literature": {
        "xiamen", "university", "campus", "factory", "business", "stock market",
        "recycling", "friction", "formula", "laboratory", "circuit", "subway", "airport",
    },
    "math": {"temple", "novel", "author", "mother", "garden", "fashion", "recycling"},
    "physics": {"temple", "novel", "author", "mother", "fashion", "recycling", "literature"},
    "business": {"square root", "temple", "novel", "friction", "formula"},
}
TERM_TRANSLATIONS = {
    "根号": "square root",
    "平方根": "square root",
    "开方": "square root",
    "二次根式": "radical expression",
    "根式": "radical expression",
    "勾股": "Pythagorean theorem",
    "直角三角形": "right triangle",
    "数轴": "number line",
    "定义域": "domain",
    "函数": "function graph",
    "曲线": "curve graph",
    "图像": "graph diagram",
    "面积": "area diagram",
    "边长": "side length",
    "正方形": "square geometry",
    "速度": "velocity physics",
    "加速度": "acceleration physics",
    "旧衣": "used clothing",
    "衣物": "clothing",
    "纺织": "textile",
    "面料": "fabric textile",
    "服装": "clothing fashion",
    "垃圾": "waste",
    "填埋": "landfill",
    "焚烧": "incineration",
    "回收": "recycling",
    "分拣": "sorting",
    "清洗": "washing",
    "消毒": "disinfection",
    "纤维": "fiber textile",
    "工厂": "factory",
    "门店": "retail store",
    "城市": "city",
    "社区": "community",
    "团队": "team",
    "消费者": "consumer",
    "品牌": "brand",
    "产品": "product",
    "材料": "material",
    "竞赛": "competition",
    "创新": "innovation",
    "创业": "entrepreneurship",
    "商业": "business",
    "市场": "market",
    "案例": "case study",
    "流程": "process workflow",
    "步骤": "process steps",
    "结构": "structure diagram",
    "实验": "experiment",
    "地图": "map",
    "语文": "Chinese literature",
    "文学": "literature",
    "课文": "textbook passage",
    "文本": "text close reading",
    "意象": "literary imagery",
    "修辞": "rhetoric figure of speech",
    "作者": "author portrait background",
    "史铁生": "Shi Tiesheng",
    "我与地坛": "I and the Temple of Earth essay",
    "地坛": "Temple of Earth Beijing",
    "历史": "history",
    "朝代": "dynasty history",
    "革命": "revolution history",
    "战争": "war history",
    "地理": "geography",
    "气候": "climate map",
    "经纬度": "latitude longitude map",
    "等高线": "contour map",
    "摩擦力": "friction force",
    "受力分析": "force diagram",
    "测力计": "spring scale experiment",
    "电路": "circuit diagram",
    "化学": "chemistry",
    "细胞": "cell biology",
    "DNA": "DNA structure",
    "算法": "algorithm diagram",
    "代码": "programming code",
}


@dataclass
class ImageQuality:
    usable: bool
    width: int = 0
    height: int = 0
    aspect_ratio: float = 0.0
    reason: str = ""


@dataclass
class ExtractedAsset:
    filename: str
    path: str
    source_type: str
    source_file: str
    slide_number: int | None
    checksum: str
    quality: ImageQuality
    semantic_hint: str = ""
    content_audit: dict[str, Any] = field(default_factory=dict)


def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat().replace("+00:00", "Z")


def file_checksum(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def clamp(value: float, low: float = 0.0, high: float = 1.0) -> float:
    return max(low, min(high, value))


def analyze_image_content(path: Path) -> dict[str, Any]:
    """Measure visual usefulness without OCR or network-dependent models."""
    if Image is None:
        return {
            "available": False,
            "accepted": True,
            "score": 0.5,
            "reason": "Pillow unavailable; content audit skipped",
        }
    try:
        with Image.open(path) as raw:
            width, height = raw.size
            aspect = width / height if height else 0.0
            img = raw.convert("RGB").resize((96, 54))
            pixels = list(img.getdata())
    except Exception as exc:
        return {
            "available": True,
            "accepted": False,
            "score": 0.0,
            "reason": f"unreadable image content: {exc}",
        }

    if not pixels:
        return {
            "available": True,
            "accepted": False,
            "score": 0.0,
            "width": 0,
            "height": 0,
            "reason": "empty pixel buffer",
        }

    grays = [(r * 299 + g * 587 + b * 114) / 1000.0 for r, g, b in pixels]
    mean = sum(grays) / len(grays)
    variance = sum((value - mean) ** 2 for value in grays) / len(grays)
    stddev = math.sqrt(variance)
    colorfulness = sum((max(r, g, b) - min(r, g, b)) / 255.0 for r, g, b in pixels) / len(pixels)

    edge_total = 0.0
    edge_hits = 0
    comparisons = 0
    grid_w, grid_h = 96, 54
    for y in range(grid_h):
        row = y * grid_w
        for x in range(grid_w):
            value = grays[row + x]
            if x + 1 < grid_w:
                diff = abs(value - grays[row + x + 1])
                edge_total += diff
                edge_hits += 1 if diff > 28 else 0
                comparisons += 1
            if y + 1 < grid_h:
                diff = abs(value - grays[row + x + grid_w])
                edge_total += diff
                edge_hits += 1 if diff > 28 else 0
                comparisons += 1

    edge_score = (edge_total / max(1, comparisons)) / 255.0
    edge_density = edge_hits / max(1, comparisons)
    bright_ratio = sum(1 for value in grays if value > 242) / len(grays)
    dark_ratio = sum(1 for value in grays if value < 18) / len(grays)
    neutral_ratio = sum(1 for r, g, b in pixels if max(r, g, b) - min(r, g, b) < 10) / len(pixels)

    tags: list[str] = []
    if stddev < 5 and (bright_ratio > 0.84 or dark_ratio > 0.84):
        tags.append("near_blank")
    if stddev < 9 and edge_score < 0.018:
        tags.append("low_detail")
    if edge_density > 0.34 and neutral_ratio > 0.62:
        tags.append("text_or_screenshot_like")
    if colorfulness < 0.035 and edge_density < 0.035:
        tags.append("flat_monochrome")

    score = 0.34
    score += clamp(stddev / 70.0) * 0.24
    score += clamp(edge_score * 6.0) * 0.22
    score += clamp(colorfulness * 2.2) * 0.16
    score += 0.04 if 0.45 <= aspect <= 2.8 else -0.08
    if "near_blank" in tags:
        score -= 0.55
    if "low_detail" in tags:
        score -= 0.30
    if "flat_monochrome" in tags:
        score -= 0.14
    score = clamp(score)

    accepted = "near_blank" not in tags and "low_detail" not in tags and score >= 0.28
    reasons = [f"{tag}" for tag in tags]
    if accepted:
        reasons.append("visual detail sufficient")
    else:
        reasons.append("below image content threshold")

    return {
        "available": True,
        "accepted": accepted,
        "score": round(score, 2),
        "width": width,
        "height": height,
        "aspect_ratio": round(aspect, 3),
        "brightness": round(mean, 2),
        "contrast": round(stddev, 2),
        "edge_score": round(edge_score, 4),
        "edge_density": round(edge_density, 4),
        "colorfulness": round(colorfulness, 4),
        "bright_ratio": round(bright_ratio, 4),
        "dark_ratio": round(dark_ratio, 4),
        "neutral_ratio": round(neutral_ratio, 4),
        "tags": tags,
        "reason": "; ".join(reasons),
    }


def audit_downloaded_visual(
    path: Path,
    source_item: dict[str, Any] | None,
    *,
    query: str = "",
    framework: dict[str, Any] | None = None,
) -> dict[str, Any]:
    metadata_audit = audit_web_source_item(source_item, query=query, framework=framework)
    content_audit = analyze_image_content(path)
    framework = framework or {}
    family = framework.get("family", "general")
    asset_mode = framework.get("asset_mode", "")

    reasons: list[str] = []
    if metadata_audit.get("reason"):
        reasons.append(f"metadata: {metadata_audit['reason']}")
    if content_audit.get("reason"):
        reasons.append(f"content: {content_audit['reason']}")

    accepted = bool(metadata_audit.get("accepted")) and bool(content_audit.get("accepted", True))
    content_tags = set(content_audit.get("tags") or [])
    if family == "humanities" and "text_or_screenshot_like" in content_tags:
        accepted = False
        reasons.append("humanities visual looks like a text screenshot, not an illustrative asset")
    if asset_mode in {"diagram", "formula_render"}:
        if content_audit.get("edge_density", 0) < 0.012 and content_audit.get("contrast", 0) < 10:
            accepted = False
            reasons.append("STEM diagram/formula visual has too little line detail")

    score = round(
        float(metadata_audit.get("score", 0.0)) * 0.62
        + float(content_audit.get("score", 0.5)) * 0.38,
        2,
    )
    return {
        "score": score,
        "accepted": accepted,
        "reason": "; ".join(part for part in reasons if part),
        "query": query or metadata_audit.get("query", ""),
        "metadata": metadata_audit,
        "content": content_audit,
        "subject": framework.get("subject") or metadata_audit.get("subject") or "general",
        "scene": framework.get("scene", "concept"),
        "asset_mode": asset_mode,
    }


def _visual_cache_key(query: str, framework: dict[str, Any] | None) -> str:
    framework = framework or {}
    basis = {
        "query": re.sub(r"\s+", " ", query or "").strip().lower(),
        "subject": framework.get("subject", "general"),
        "scene": framework.get("scene", "concept"),
        "asset_mode": framework.get("asset_mode", ""),
        "purpose": framework.get("cache_purpose", ""),
        "variant": framework.get("background_variant", ""),
    }
    return hashlib.sha256(json.dumps(basis, sort_keys=True, ensure_ascii=False).encode("utf-8")).hexdigest()[:20]


def _write_image_source_item(output_dir: Path, item: dict[str, Any]) -> None:
    manifest_path = output_dir / "image_sources.json"
    try:
        payload = json.loads(manifest_path.read_text(encoding="utf-8")) if manifest_path.exists() else {}
    except (OSError, json.JSONDecodeError):
        payload = {}
    items = [entry for entry in payload.get("items", []) if entry.get("filename") != item.get("filename")]
    items.append(item)
    payload["generated_at"] = now_iso()
    payload["items"] = items
    payload.setdefault("license_verification", "provider metadata used; manual review recommended for external delivery")
    manifest_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def restore_visual_cache(
    query: str,
    framework: dict[str, Any] | None,
    output_path: Path,
    output_dir: Path,
) -> dict[str, Any] | None:
    key = _visual_cache_key(query, framework)
    meta_path = VISUAL_CACHE_DIR / f"{key}.json"
    if not meta_path.exists():
        return None
    try:
        meta = json.loads(meta_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None
    cached_path = Path(meta.get("path", ""))
    if not cached_path.exists():
        return None
    output_dir.mkdir(parents=True, exist_ok=True)
    shutil.copy2(cached_path, output_path)
    item = dict(meta.get("source_item") or {})
    item.update(
        {
            "filename": output_path.name,
            "search_query": query,
            "status": "cache_hit",
            "cache_key": key,
            "cache_path": str(cached_path),
        }
    )
    _write_image_source_item(output_dir, item)
    return {"cache_key": key, "cache_path": str(cached_path), "source_item": item}


def store_visual_cache(
    query: str,
    framework: dict[str, Any] | None,
    image_path: Path,
    source_item: dict[str, Any] | None,
    audit: dict[str, Any],
) -> dict[str, Any]:
    key = _visual_cache_key(query, framework)
    VISUAL_CACHE_DIR.mkdir(parents=True, exist_ok=True)
    suffix = image_path.suffix.lower() or ".jpg"
    cache_path = VISUAL_CACHE_DIR / f"{key}{suffix}"
    shutil.copy2(image_path, cache_path)
    meta = {
        "cache_key": key,
        "created_at": now_iso(),
        "query": query,
        "framework": {
            "subject": (framework or {}).get("subject", "general"),
            "family": (framework or {}).get("family", "general"),
            "scene": (framework or {}).get("scene", "concept"),
            "asset_mode": (framework or {}).get("asset_mode", ""),
        },
        "path": str(cache_path),
        "source_item": source_item or {},
        "audit": audit,
    }
    (VISUAL_CACHE_DIR / f"{key}.json").write_text(
        json.dumps(meta, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    return {"cache_key": key, "cache_path": str(cache_path)}


def sanitize_token(value: str, fallback: str = "asset") -> str:
    safe = re.sub(r"[^0-9A-Za-z_-]+", "_", value.strip())
    safe = re.sub(r"_+", "_", safe).strip("_")
    return safe[:80] or fallback


def normalize_text(value: str) -> str:
    return re.sub(r"\s+", " ", str(value or "")).strip().lower()


def extract_content_terms(text: str, *, max_terms: int = 12) -> list[str]:
    """Extract compact topic terms from Chinese/English slide text."""
    normalized = normalize_text(text)
    terms: list[str] = []

    known_terms = sorted(
        set(CONCRETE_KEYWORDS + ABSTRACT_KEYWORDS + tuple(TERM_TRANSLATIONS)),
        key=len,
        reverse=True,
    )
    for term in known_terms:
        if not term or len(term) < 2:
            continue
        if term.lower() in normalized and term not in CHINESE_STOPWORDS:
            terms.append(term)

    for word in re.findall(r"[A-Za-z][A-Za-z0-9+-]{2,}", text):
        lower = word.lower()
        if lower not in ENGLISH_STOPWORDS and lower not in GENERIC_ASSET_TOKENS:
            terms.append(lower)

    for phrase in re.findall(r"[\u4e00-\u9fff]{2,8}", text):
        if phrase in CHINESE_STOPWORDS:
            continue
        if len(phrase) <= 2:
            terms.append(phrase)
            continue
        # Long Chinese runs often contain punctuation-stripped sentences. Keep
        # informative 2-4 char windows that match known domain terms first.
        for width in (4, 3, 2):
            for i in range(0, max(0, len(phrase) - width + 1)):
                piece = phrase[i : i + width]
                if piece in TERM_TRANSLATIONS or piece in CONCRETE_KEYWORDS or piece in ABSTRACT_KEYWORDS:
                    terms.append(piece)

    unique: list[str] = []
    seen: set[str] = set()
    for term in terms:
        key = term.lower()
        if key in seen:
            continue
        seen.add(key)
        unique.append(term)
        if len(unique) >= max_terms:
            break
    return unique


def asset_hint_text(asset: "ExtractedAsset") -> str:
    pieces = [
        asset.filename,
        Path(asset.path).stem,
        Path(asset.source_file).stem,
        asset.source_type,
        asset.semantic_hint,
    ]
    return " ".join(piece for piece in pieces if piece)


def extract_asset_terms(asset: "ExtractedAsset") -> list[str]:
    hint = asset_hint_text(asset)
    terms: list[str] = []
    for word in re.findall(r"[A-Za-z][A-Za-z0-9+-]{2,}", hint):
        lower = word.lower()
        if lower not in GENERIC_ASSET_TOKENS and lower not in ENGLISH_STOPWORDS:
            terms.append(lower)
    for phrase in re.findall(r"[\u4e00-\u9fff]{2,8}", hint):
        if phrase not in CHINESE_STOPWORDS:
            terms.append(phrase)
    unique: list[str] = []
    seen: set[str] = set()
    for term in terms:
        key = term.lower()
        if key not in seen:
            seen.add(key)
            unique.append(term)
    return unique


def looks_generic_asset_name(asset: "ExtractedAsset") -> bool:
    hint = " ".join([asset.filename, Path(asset.path).stem, asset.source_type]).lower()
    if re.search(r"(ppt_)?slide[_-]?\d+[_-]?image[_-]?\d+", hint):
        return True
    if re.search(r"ppt[_-]?media[_-]?\d+", hint):
        return True
    semantic_terms = extract_asset_terms(asset)
    return len(semantic_terms) == 0


def score_asset_relevance(text: str, asset: "ExtractedAsset") -> dict[str, Any]:
    content_terms = extract_content_terms(text, max_terms=16)
    asset_terms = extract_asset_terms(asset)
    content_lowers = {term.lower() for term in content_terms}
    asset_lowers = {term.lower() for term in asset_terms}
    overlap = sorted(content_lowers & asset_lowers)

    score = 0.0
    reasons: list[str] = []
    if not asset.quality.usable:
        reasons.append(f"quality rejected: {asset.quality.reason or 'not usable'}")
        return {
            "score": 0.0,
            "accepted": False,
            "reason": "; ".join(reasons),
            "content_terms": content_terms,
            "asset_terms": asset_terms,
            "overlap_terms": overlap,
            "generic_name": looks_generic_asset_name(asset),
        }

    if asset.slide_number is not None:
        score += 0.12
        reasons.append("same-slide source image")
    if overlap:
        score += min(0.58, 0.24 + 0.12 * len(overlap))
        reasons.append("text/asset term overlap: " + ", ".join(overlap[:5]))
    if asset.semantic_hint:
        score += 0.08
        reasons.append("has source alt/name hint")
    if looks_generic_asset_name(asset):
        score -= 0.18
        reasons.append("generic PPT media filename; no semantic proof")
    if asset.source_type == "source_markdown_image" and asset_terms:
        score += 0.08
        reasons.append("markdown image path carries topic terms")

    score = max(0.0, min(1.0, score))
    accepted = score >= RELEVANCE_ACCEPT_THRESHOLD
    if not accepted and not reasons:
        reasons.append("no semantic overlap with slide content")
    elif not accepted:
        reasons.append("below relevance threshold")

    return {
        "score": round(score, 2),
        "accepted": accepted,
        "reason": "; ".join(reasons),
        "content_terms": content_terms,
        "asset_terms": asset_terms,
        "overlap_terms": overlap,
        "generic_name": looks_generic_asset_name(asset),
    }


def _web_text_terms(text: str, *, max_terms: int = 24) -> list[str]:
    """Extract comparable terms from web image metadata and search queries."""
    if not text:
        return []
    text = re.sub(r"https?://", " ", str(text))
    text = re.sub(r"[_/\\?#=&.%+-]+", " ", text)
    terms: list[str] = []
    for word in re.findall(r"[A-Za-z][A-Za-z0-9-]{2,}", text):
        lower = word.lower().strip("-")
        if not lower or lower in ENGLISH_STOPWORDS or lower in GENERIC_ASSET_TOKENS:
            continue
        terms.append(lower)
    for phrase in re.findall(r"[\u4e00-\u9fff]{2,8}", text):
        if phrase not in CHINESE_STOPWORDS:
            terms.append(phrase)
    unique: list[str] = []
    seen: set[str] = set()
    for term in terms:
        key = term.lower()
        if key in seen:
            continue
        seen.add(key)
        unique.append(term)
        if len(unique) >= max_terms:
            break
    return unique


def _framework_distinctive_terms(framework: dict[str, Any] | None) -> list[str]:
    if not framework:
        return []
    terms: list[str] = []
    for key in ("subject", "subject_label", "scene"):
        value = framework.get(key)
        if value and value != "general":
            terms.extend(_web_text_terms(str(value), max_terms=4))
    for hint in framework.get("query_hints") or []:
        terms.extend(_web_text_terms(str(hint), max_terms=6))
    return list(dict.fromkeys(terms))


def audit_web_source_item(
    source_item: dict[str, Any] | None,
    *,
    query: str = "",
    framework: dict[str, Any] | None = None,
    min_score: float = WEB_RELEVANCE_ACCEPT_THRESHOLD,
) -> dict[str, Any]:
    """Audit a downloaded web image before it is allowed into rendering.

    Provider scoring already filters obvious misses.  This second pass is
    stricter for our video pipeline: it records why a visual was accepted and
    rejects assets whose metadata only matches generic words.
    """
    source_item = source_item or {}
    query = query or str(source_item.get("search_query") or "")
    metadata = " ".join(
        str(source_item.get(key) or "")
        for key in ("title", "author", "source_page_url", "download_url", "provider")
    )
    query_terms = _web_text_terms(query, max_terms=18)
    metadata_terms = _web_text_terms(metadata, max_terms=32)
    framework_terms = _framework_distinctive_terms(framework)

    metadata_lowers = {term.lower() for term in metadata_terms}
    query_lowers = {term.lower() for term in query_terms}
    framework_lowers = {term.lower() for term in framework_terms}
    overlap = sorted((query_lowers | framework_lowers) & metadata_lowers)

    score = 0.0
    reasons: list[str] = []
    if overlap:
        score += min(0.70, 0.20 + 0.10 * len(overlap))
        reasons.append("metadata overlap: " + ", ".join(overlap[:6]))
    if source_item.get("width") and source_item.get("height"):
        pixels = int(source_item.get("width") or 0) * int(source_item.get("height") or 0)
        if pixels >= 900_000:
            score += 0.10
            reasons.append("sufficient downloaded resolution")
    if source_item.get("provider") in {"pexels", "pixabay"}:
        score += 0.04
        reasons.append("curated stock provider")
    if source_item.get("license_tier") in {"no-attribution", "attribution-required"}:
        score += 0.04
        reasons.append("usable license tier")

    subject = (framework or {}).get("subject", "")
    mismatch_terms = WEB_SUBJECT_MISMATCH_TERMS.get(subject, set())
    metadata_lower_text = metadata.lower()
    mismatches = sorted(term for term in mismatch_terms if term in metadata_lower_text and term not in query.lower())
    if mismatches:
        score -= min(0.32, 0.12 * len(mismatches))
        reasons.append("subject mismatch hints: " + ", ".join(mismatches[:4]))

    if query_terms and not overlap:
        reasons.append("no search-query terms found in image metadata")
    if not metadata_terms:
        reasons.append("no usable metadata to verify relevance")

    score = max(0.0, min(1.0, score))
    accepted = score >= min_score
    if not accepted:
        reasons.append("below web relevance threshold")

    return {
        "score": round(score, 2),
        "accepted": accepted,
        "threshold": min_score,
        "reason": "; ".join(reasons),
        "query": query,
        "query_terms": query_terms,
        "metadata_terms": metadata_terms,
        "framework_terms": framework_terms,
        "overlap_terms": overlap,
        "subject": subject or "general",
    }


def _read_image_sources(output_dir: Path) -> dict[str, dict[str, Any]]:
    path = output_dir / "image_sources.json"
    if not path.exists():
        return {}
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}
    return {
        str(item.get("filename") or ""): item
        for item in payload.get("items", [])
        if isinstance(item, dict) and item.get("filename")
    }


def _quarantine_rejected_asset(path: Path) -> str:
    rejected_dir = path.parent / "rejected"
    rejected_dir.mkdir(parents=True, exist_ok=True)
    dst = rejected_dir / path.name
    if dst.exists():
        dst = rejected_dir / f"{path.stem}_{int(datetime.now().timestamp())}{path.suffix}"
    try:
        shutil.move(str(path), str(dst))
        return str(dst)
    except OSError:
        return str(path)


def is_no_content_placeholder(text: str) -> bool:
    cleaned = clean_line(text).strip("_ ").lower().rstrip(".")
    return "no extractable text content" in cleaned


def validate_image(path: Path, min_width: int, min_height: int) -> ImageQuality:
    if Image is None:
        probed = probe_image_with_ffprobe(path)
        if probed:
            width, height = probed
            aspect = width / height if height else 0.0
            if width < min_width or height < min_height:
                return ImageQuality(False, width, height, aspect, "too small")
            if aspect < 0.25 or aspect > 4.0:
                return ImageQuality(False, width, height, aspect, "extreme aspect ratio")
            return ImageQuality(True, width, height, aspect, "ok via ffprobe")
        if path.stat().st_size < 80_000:
            return ImageQuality(False, reason="Pillow unavailable and file too small to trust")
        return ImageQuality(True, reason="Pillow unavailable; accepted by file size only")
    try:
        with Image.open(path) as img:
            width, height = img.size
            aspect = width / height if height else 0.0
            if width < min_width or height < min_height:
                return ImageQuality(False, width, height, aspect, "too small")
            if aspect < 0.25 or aspect > 4.0:
                return ImageQuality(False, width, height, aspect, "extreme aspect ratio")
            if ImageStat is not None:
                stat = ImageStat.Stat(img.convert("L").resize((64, 64)))
                if stat.var and stat.var[0] < 3:
                    return ImageQuality(False, width, height, aspect, "near-blank image")
            content_audit = analyze_image_content(path)
            if content_audit.get("accepted") is False:
                return ImageQuality(
                    False,
                    width,
                    height,
                    aspect,
                    f"content audit rejected: {content_audit.get('reason', '')}",
                )
            return ImageQuality(True, width, height, aspect, "ok")
    except Exception as exc:
        return ImageQuality(False, reason=f"unreadable: {exc}")


def probe_image_with_ffprobe(path: Path) -> tuple[int, int] | None:
    """Return image dimensions using ffprobe when Pillow is unavailable."""
    try:
        result = subprocess.run(
            [
                "ffprobe",
                "-v",
                "error",
                "-select_streams",
                "v:0",
                "-show_entries",
                "stream=width,height",
                "-of",
                "json",
                str(path),
            ],
            check=True,
            capture_output=True,
            text=True,
        )
        streams = json.loads(result.stdout).get("streams") or []
        if not streams:
            return None
        width = int(streams[0].get("width") or 0)
        height = int(streams[0].get("height") or 0)
        if width and height:
            return width, height
    except Exception:
        return None
    return None


def copy_unique_image(
    src: Path,
    dest_dir: Path,
    *,
    filename_hint: str,
    seen: dict[str, ExtractedAsset],
    source_type: str,
    source_file: Path,
    slide_number: int | None,
    min_width: int,
    min_height: int,
    semantic_hint: str = "",
) -> ExtractedAsset | None:
    if not src.exists() or src.suffix.lower() not in IMAGE_EXTENSIONS:
        return None
    checksum = file_checksum(src)
    if checksum in seen:
        return None

    quality = validate_image(src, min_width, min_height)
    content_audit = analyze_image_content(src)
    suffix = src.suffix.lower()
    dest_dir.mkdir(parents=True, exist_ok=True)
    base = sanitize_token(filename_hint)
    dst = dest_dir / f"{base}{suffix}"
    if dst.exists() and file_checksum(dst) == checksum:
        asset = ExtractedAsset(
            filename=dst.name,
            path=str(dst),
            source_type=source_type,
            source_file=str(source_file),
            slide_number=slide_number,
            checksum=checksum,
            quality=quality,
            semantic_hint=semantic_hint,
            content_audit=content_audit,
        )
        seen[checksum] = asset
        return asset
    counter = 2
    while dst.exists():
        dst = dest_dir / f"{base}_{counter}{suffix}"
        counter += 1

    shutil.copy2(src, dst)
    asset = ExtractedAsset(
        filename=dst.name,
        path=str(dst),
        source_type=source_type,
        source_file=str(source_file),
        slide_number=slide_number,
        checksum=checksum,
        quality=quality,
        semantic_hint=semantic_hint,
        content_audit=content_audit,
    )
    seen[checksum] = asset
    return asset


def iter_markdown_image_refs(markdown_path: Path) -> Iterable[tuple[int | None, Path, str]]:
    slide_number: int | None = None
    image_re = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")
    for line in markdown_path.read_text(encoding="utf-8", errors="ignore").splitlines():
        match = re.match(r"^##\s+Slide\s+(\d+)", line.strip(), re.I)
        if match:
            slide_number = int(match.group(1))
            continue
        for image_match in image_re.finditer(line):
            rel = image_match.group(1).strip()
            if rel.startswith(("http://", "https://")):
                continue
            alt_match = re.search(r"!\[([^\]]*)\]\(" + re.escape(rel) + r"\)", line)
            alt_text = alt_match.group(1).strip() if alt_match else ""
            yield slide_number, (markdown_path.parent / rel).resolve(), alt_text


def extract_images_from_markdown_sources(
    project: Path,
    dest_dir: Path,
    seen: dict[str, ExtractedAsset],
    min_width: int,
    min_height: int,
) -> list[ExtractedAsset]:
    assets: list[ExtractedAsset] = []
    for md_path in sorted((project / "sources").glob("*.md")):
        for slide_number, image_path, alt_text in iter_markdown_image_refs(md_path):
            hint = f"src_slide_{slide_number or 0:02d}_{image_path.stem}"
            asset = copy_unique_image(
                image_path,
                dest_dir,
                filename_hint=hint,
                seen=seen,
                source_type="source_markdown_image",
                source_file=md_path,
                slide_number=slide_number,
                min_width=min_width,
                min_height=min_height,
                semantic_hint=alt_text,
            )
            if asset:
                assets.append(asset)
    return assets


def extract_pptx_with_python_pptx(
    pptx_path: Path,
    dest_dir: Path,
    seen: dict[str, ExtractedAsset],
    min_width: int,
    min_height: int,
) -> list[ExtractedAsset]:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except Exception:
        return []

    assets: list[ExtractedAsset] = []
    prs = Presentation(str(pptx_path))

    def walk_shapes(shapes: Any) -> Iterable[Any]:
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk_shapes(shape.shapes)
            else:
                yield shape

    def shape_semantic_hint(shape: Any) -> str:
        pieces: list[str] = []
        for attr in ("name", "alternative_text", "alt_text"):
            value = getattr(shape, attr, "")
            if value:
                pieces.append(str(value))
        try:
            for node in shape._element.xpath(".//p:cNvPr"):
                for attr in ("name", "descr", "title"):
                    value = node.get(attr)
                    if value:
                        pieces.append(str(value))
        except Exception:
            pass
        unique: list[str] = []
        seen_piece: set[str] = set()
        for piece in pieces:
            cleaned = re.sub(r"\s+", " ", piece).strip()
            if cleaned and cleaned not in seen_piece:
                seen_piece.add(cleaned)
                unique.append(cleaned)
        return " ".join(unique)

    for slide_idx, slide in enumerate(prs.slides, 1):
        image_idx = 0
        for shape in walk_shapes(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                image = shape.image
                image_idx += 1
                suffix = f".{(image.ext or 'png').lower()}"
                tmp = dest_dir / f"__tmp_slide_{slide_idx:02d}_{image_idx:02d}{suffix}"
                tmp.write_bytes(image.blob)
                asset = copy_unique_image(
                    tmp,
                    dest_dir,
                    filename_hint=f"ppt_slide_{slide_idx:02d}_image_{image_idx:02d}",
                    seen=seen,
                    source_type="source_pptx_picture",
                    source_file=pptx_path,
                    slide_number=slide_idx,
                    min_width=min_width,
                    min_height=min_height,
                    semantic_hint=shape_semantic_hint(shape),
                )
                tmp.unlink(missing_ok=True)
                if asset:
                    assets.append(asset)
            except Exception:
                continue
    return assets


def extract_pptx_media_zip(
    pptx_path: Path,
    dest_dir: Path,
    seen: dict[str, ExtractedAsset],
    min_width: int,
    min_height: int,
) -> list[ExtractedAsset]:
    assets: list[ExtractedAsset] = []
    try:
        with zipfile.ZipFile(pptx_path) as zf:
            names = [name for name in zf.namelist() if name.startswith("ppt/media/")]
            for idx, name in enumerate(names, 1):
                suffix = Path(name).suffix.lower()
                if suffix not in IMAGE_EXTENSIONS:
                    continue
                tmp = dest_dir / f"__tmp_media_{idx:03d}{suffix}"
                tmp.write_bytes(zf.read(name))
                asset = copy_unique_image(
                    tmp,
                    dest_dir,
                    filename_hint=f"ppt_media_{idx:03d}",
                    seen=seen,
                    source_type="source_pptx_media",
                    source_file=pptx_path,
                    slide_number=None,
                    min_width=min_width,
                    min_height=min_height,
                )
                tmp.unlink(missing_ok=True)
                if asset:
                    assets.append(asset)
    except zipfile.BadZipFile:
        return []
    return assets


def extract_source_images(
    project: Path,
    min_width: int,
    min_height: int,
) -> list[ExtractedAsset]:
    dest_dir = project / "images" / "extracted"
    dest_dir.mkdir(parents=True, exist_ok=True)
    seen: dict[str, ExtractedAsset] = {}
    assets: list[ExtractedAsset] = []
    assets.extend(extract_images_from_markdown_sources(project, dest_dir, seen, min_width, min_height))

    for pptx_path in sorted((project / "sources").iterdir() if (project / "sources").exists() else []):
        if pptx_path.suffix.lower() not in PPT_EXTENSIONS:
            continue
        extracted = extract_pptx_with_python_pptx(pptx_path, dest_dir, seen, min_width, min_height)
        if not extracted:
            extracted = extract_pptx_media_zip(pptx_path, dest_dir, seen, min_width, min_height)
        assets.extend(extracted)
    return assets


def clean_line(line: str) -> str:
    line = re.sub(r"\s+", " ", line.strip())
    line = re.sub(r"!\[[^\]]*\]\([^)]+\)", "", line).strip()
    return line


def load_slide_texts(project: Path) -> list[dict[str, Any]]:
    source_slides: dict[int, list[str]] = {}
    for md_path in sorted((project / "sources").glob("*.md")):
        current: int | None = None
        for raw in md_path.read_text(encoding="utf-8", errors="ignore").splitlines():
            match = re.match(r"^##\s+Slide\s+(\d+)", raw.strip(), re.I)
            if match:
                current = int(match.group(1))
                source_slides.setdefault(current, [])
                continue
            if current is None:
                continue
            if raw.startswith("## ") or raw.startswith("### Speaker Notes"):
                current = None
                continue
            line = clean_line(raw)
            if not line or is_no_content_placeholder(line) or line.startswith(("#", "-", "Source:", "Total slides:")):
                continue
            source_slides[current].append(line)

    if source_slides:
        return [
            {
                "slide_number": slide_num,
                "title": lines[0] if lines else f"Slide {slide_num}",
                "segments": lines[1:] or lines[:1],
            }
            for slide_num, lines in sorted(source_slides.items())
        ]

    structure_file = project / "slide_structure.json"
    if not structure_file.exists():
        return []
    data = json.loads(structure_file.read_text(encoding="utf-8"))
    slides = data.get("slides", data) if isinstance(data, dict) else data
    result = []
    for slide in slides:
        title = slide.get("title") or f"Slide {slide.get('slide_number')}"
        segments = []
        segments.extend(slide.get("bullets") or [])
        segments.extend(slide.get("paragraphs") or [])
        result.append(
            {
                "slide_number": int(slide.get("slide_number") or len(result) + 1),
                "title": title,
                "segments": [clean_line(s) for s in segments if clean_line(s)],
            }
        )
    return result


def segment_subject_framework(title: str, segment: str, slide_num: int) -> dict[str, Any]:
    if framework_for_slide is None:
        return {}
    return framework_for_slide(title, [segment], slide_num)


def score_segment_need(text: str, framework: dict[str, Any] | None = None) -> float:
    if not text:
        return 0.0
    cleaned = re.sub(r"\s+", "", text).strip()
    if cleaned in {"目录", "感谢观看", "谢谢观看", "thankyou", "thanks"}:
        return 0.0
    if re.fullmatch(r"\d{1,2}", cleaned):
        return 0.0
    if re.search(r"(目录|agenda)", text, re.I) and len(cleaned) < 18:
        return 0.0
    score = 0.0
    if any(key in text for key in CONCRETE_KEYWORDS):
        score += 0.45
    if any(key in text for key in ABSTRACT_KEYWORDS):
        score += 0.18
    if extract_content_terms(text, max_terms=4):
        score += 0.2
    if any(key in text for key in NO_IMAGE_KEYWORDS):
        score -= 0.28
    if re.search(r"\d", text) and len(text) < 50:
        score -= 0.15
    if len(text) >= 28:
        score += 0.12
    if framework:
        policy = framework.get("image_policy", "optional")
        asset_mode = framework.get("asset_mode", "none")
        family = framework.get("family", "general")
        if policy == "required":
            score += 0.55
        elif policy == "prefer":
            score += 0.38
        elif policy == "optional":
            score += 0.10
        elif policy == "none":
            score -= 0.30
        if asset_mode in {"search", "diagram", "formula_render"}:
            score += 0.20
        elif asset_mode == "generate":
            score += 0.14
        if family == "humanities" and policy != "none":
            score += 0.12
    return max(0.0, min(1.0, score))


def select_relevant_existing_asset(text: str, slide_assets: list[ExtractedAsset]) -> tuple[ExtractedAsset | None, dict[str, Any]]:
    scored = []
    for asset in slide_assets:
        relevance = score_asset_relevance(text, asset)
        scored.append({"asset": asset, "relevance": relevance})
    scored.sort(key=lambda item: item["relevance"]["score"], reverse=True)
    accepted = next((item for item in scored if item["relevance"]["accepted"]), None)
    diagnostics = {
        "threshold": RELEVANCE_ACCEPT_THRESHOLD,
        "candidates": [
            {
                "filename": item["asset"].filename,
                "path": item["asset"].path,
                "score": item["relevance"]["score"],
                "accepted": item["relevance"]["accepted"],
                "reason": item["relevance"]["reason"],
                "content_terms": item["relevance"]["content_terms"],
                "asset_terms": item["relevance"]["asset_terms"],
                "overlap_terms": item["relevance"]["overlap_terms"],
                "generic_name": item["relevance"]["generic_name"],
            }
            for item in scored
        ],
    }
    if accepted:
        return accepted["asset"], diagnostics
    return None, diagnostics


def choose_mode(
    text: str,
    slide_assets: list[ExtractedAsset],
    framework: dict[str, Any] | None = None,
) -> tuple[str, ExtractedAsset | None, dict[str, Any]]:
    selected_asset, diagnostics = select_relevant_existing_asset(text, slide_assets)
    if selected_asset:
        return "existing", selected_asset, diagnostics
    if framework:
        policy = framework.get("image_policy", "optional")
        asset_mode = framework.get("asset_mode", "none")
        if policy == "none" or asset_mode == "none":
            return "none", None, diagnostics
        if asset_mode in {"search", "diagram", "formula_render"}:
            return "search", None, diagnostics
        if asset_mode == "generate":
            return "generate", None, diagnostics
    if any(key in text for key in CONCRETE_KEYWORDS):
        return "search", None, diagnostics
    if extract_content_terms(text, max_terms=3):
        return "search", None, diagnostics
    if any(key in text for key in ABSTRACT_KEYWORDS):
        return "generate", None, diagnostics
    return "none", None, diagnostics


def build_search_query(text: str, framework: dict[str, Any] | None = None) -> str:
    if framework:
        hints = [hint for hint in framework.get("query_hints", []) if hint]
        if hints:
            return hints[0]
    matches = [
        (len(keys), sum(len(key) for key in keys), query)
        for keys, query in SEARCH_QUERY_RULES
        if all(key in text for key in keys)
    ]
    if matches:
        matches.sort(reverse=True)
        return matches[0][2]
    for keys, query in SEARCH_QUERY_RULES:
        if any(key in text for key in keys):
            return query
    terms = extract_content_terms(text, max_terms=5)
    translated: list[str] = []
    for term in terms:
        translated.extend((TERM_TRANSLATIONS.get(term) or term).split())
    cleaned: list[str] = []
    seen: set[str] = set()
    for term in translated:
        lower = term.lower()
        if len(lower) < 2 or lower in ENGLISH_STOPWORDS or lower in seen:
            continue
        seen.add(lower)
        cleaned.append(term)
    if cleaned:
        suffix = "diagram" if any(key in text for key in ("图", "示意", "结构", "流程", "步骤", "函数", "数轴", "例题")) else "photo"
        return " ".join(cleaned[:5] + [suffix])
    latin_words = re.findall(r"[A-Za-z][A-Za-z-]{2,}", text)
    if latin_words:
        return " ".join(latin_words[:4])
    if any(key in text for key in ("环保", "可持续", "绿色")):
        return "sustainable fashion"
    if any(key in text for key in ("团队", "成员")):
        return "business team meeting"
    chinese_terms = re.findall(r"[\u4e00-\u9fff]{2,4}", text)
    if chinese_terms:
        return " ".join(chinese_terms[:4])
    return "presentation topic visual"


def should_plan_background_decor(framework: dict[str, Any] | None, text: str) -> bool:
    framework = framework or {}
    if framework.get("image_policy") == "none":
        return False
    family = framework.get("family", "general")
    subject = framework.get("subject", "general")
    scene = framework.get("scene", "concept")
    if family == "humanities":
        return True
    if family == "stem" and scene in {"cover", "definition", "formula", "experiment", "example", "case"}:
        return True
    if subject in {"geography", "history", "business"}:
        return True
    return bool(extract_content_terms(text, max_terms=3))


def build_background_decor_query(text: str, framework: dict[str, Any] | None = None) -> str:
    framework = framework or {}
    subject = framework.get("subject", "general")
    family = framework.get("family", "general")
    scene = framework.get("scene", "concept")
    motifs = set(framework.get("content_motifs") or [])
    query_hints = [str(hint) for hint in framework.get("query_hints") or [] if hint]
    base = query_hints[0] if query_hints else build_search_query(text, framework)

    if subject == "literature":
        temple_context = "Temple of Earth" in base or "Shi Tiesheng" in base or "temple_gate" in motifs
        if temple_context:
            scene_queries = {
                "cover": "Temple of Earth Beijing old cypress gate quiet garden path literary documentary photo",
                "objective": "Chinese literature close reading open book warm paper desk photo",
                "reading_path": "Beijing park stone path old map route paper texture literary lesson",
                "quote_analysis": "old Chinese book page handwritten notes warm paper texture photo",
                "close_reading": "literary annotation notebook book margin warm desk photo",
                "character_analysis": "mother memory park path back view warm documentary photo",
                "emotion_curve": "autumn park path fallen leaves quiet life reflection photo",
                "comparison_reading": "two open books side by side literature notes paper texture photo",
                "discussion": "classroom literature discussion notebooks warm desk photo",
                "writing_task": "fountain pen notebook literary writing warm paper photo",
                "summary": "Temple of Earth Beijing quiet autumn garden path closing photo",
            }
            if "mother_shadow" in motifs and scene in {"cover", "character_analysis"}:
                return "mother memory quiet park path warm literary documentary photo"
            if "quote_scroll" in motifs and scene in {"quote_analysis", "close_reading"}:
                return "Chinese literature book page annotation warm paper texture photo"
            if "autumn_path" in motifs or scene == "emotion_curve":
                return "autumn park path fallen leaves warm literary background photo"
            return scene_queries.get(
                str(scene),
                "Temple of Earth Beijing quiet garden path cypress literary documentary photo",
            )
        if "sea_waves" in motifs or any(key in text for key in ("观沧海", "沧海", "洪波", "山岛")):
            return "Chinese ink sea waves mountain paper texture background"
        if "autumn_path" in motifs or any(key in text for key in ("秋天", "落叶", "秋风", "园路")):
            return "autumn leaves garden path warm paper texture background"
        scene_queries = {
            "cover": f"{base} literary documentary photo background",
            "objective": "open book literature notes warm paper desk photo",
            "reading_path": "literature reading path notebook timeline paper texture photo",
            "quote_analysis": "book page quotation annotation warm paper texture photo",
            "close_reading": "close reading annotation notebook warm desk photo",
            "character_analysis": "literary character memory portrait silhouette background photo",
            "emotion_curve": "quiet path fallen leaves emotional reflection background photo",
            "discussion": "classroom literature discussion notebooks warm desk photo",
            "writing_task": "fountain pen notebook literature writing warm paper photo",
            "summary": "open book warm sunlight paper texture closing background photo",
        }
        return scene_queries.get(str(scene), f"{base} literary paper texture background")
    if subject == "history":
        return f"{base} archive paper map texture background"
    if subject == "geography":
        return f"{base} atlas map contour texture background"
    if subject == "math":
        return f"{base} graph paper formula texture background"
    if subject == "physics":
        return f"{base} physics experiment blackboard texture background"
    if subject == "chemistry":
        return f"{base} chemistry lab glass texture background"
    if subject == "biology":
        return f"{base} biology microscope natural texture background"
    if subject == "computer":
        return f"{base} circuit board data texture background"
    if family == "business":
        return f"{base} business dashboard subtle background"
    return f"{base} subtle educational background texture"


def build_generation_prompt(text: str, title: str, framework: dict[str, Any] | None = None) -> str:
    compact = re.sub(r"\s+", " ", f"{title}. {text}").strip()
    if framework:
        family = framework.get("family", "general")
        subject = framework.get("subject_label", framework.get("subject", ""))
        scene = framework.get("scene", "concept")
        if family == "humanities":
            style = (
                "Create a literary editorial illustration or documentary-style image. "
                "Use warm paper texture, real contextual details, and a restrained classroom aesthetic. "
            )
        elif family == "stem":
            style = (
                "Create a clean educational diagram or experiment-style visual. "
                "Make the concept readable, precise, and suitable for a science classroom. "
            )
        elif family == "business":
            style = "Create a realistic business case visual with clear evidence and no decorative clutter. "
        else:
            style = "Create a clean educational presentation visual. "
        return (
            f"{style}No text, no logos, no watermarks. Subject area: {subject}; scene: {scene}. "
            f"Topic: {compact[:260]}"
        )
    return (
        "Create a clean, realistic editorial image for a business presentation. "
        "No text, no logos, no watermarks. Subject: "
        f"{compact[:260]}"
    )


def visual_decision_reason(
    actionable: dict[str, Any] | None,
    framework: dict[str, Any],
    rejection_summary: str,
) -> str:
    if actionable and actionable["mode"] == "existing":
        return "reuse source image after relevance check"
    if actionable and actionable["mode"] == "search":
        asset_mode = framework.get("asset_mode")
        subject = framework.get("subject")
        if subject == "geography":
            return "geography scene, search map/atlas diagram matched to the topic"
        if subject == "history":
            return "history scene, search timeline/map/source-image evidence"
        if asset_mode == "formula_render":
            return "STEM formula/exercise scene, search formula or worked-example diagram"
        if asset_mode == "diagram":
            return "STEM concept scene, search precise diagram or experiment visual"
        if framework.get("family") == "humanities":
            return "humanities scene, search contextual illustration/photo instead of plain text"
        return "concrete subject, use licensed search"
    if actionable and actionable["mode"] == "generate":
        if framework.get("family") == "humanities":
            return "humanities scene, generate literary/editorial supporting image"
        return "abstract/custom concept, use generation"
    return rejection_summary or "text/data/structure can be rendered without an image"


def summarize_asset_rejections(segment_decisions: list[dict[str, Any]], slide_assets: list[ExtractedAsset]) -> str:
    if not slide_assets:
        return ""
    best: dict[str, Any] | None = None
    for segment in segment_decisions:
        for candidate in segment.get("existing_asset_relevance", {}).get("candidates", []):
            if best is None or float(candidate.get("score") or 0.0) > float(best.get("score") or 0.0):
                best = candidate
    if best:
        return (
            "source image rejected by relevance check "
            f"(best={best.get('filename')}, score={best.get('score')}, reason={best.get('reason')})"
        )
    if any(asset.quality.usable for asset in slide_assets):
        return "source images available but no segment needs an image strongly enough"
    return "source images failed quality checks"


def plan_visual_assets(slides: list[dict[str, Any]], assets: list[ExtractedAsset], project_context: str = "") -> list[dict[str, Any]]:
    by_slide: dict[int, list[ExtractedAsset]] = {}
    for asset in assets:
        if asset.slide_number is not None:
            by_slide.setdefault(asset.slide_number, []).append(asset)

    decisions: list[dict[str, Any]] = []
    for slide in slides:
        slide_num = int(slide["slide_number"])
        title = slide.get("title", "")
        segments = slide.get("segments") or [title]
        slide_assets = by_slide.get(slide_num, [])
        deck_context = clean_line(project_context)
        slide_framework = segment_subject_framework(title, f"{deck_context} {' '.join(segments[:4])}", slide_num)

        segment_decisions = []
        for idx, segment in enumerate(segments[:8], 1):
            segment_text = f"{title} {segment}"
            framework = segment_subject_framework(title, f"{deck_context} {segment}", slide_num)
            score = score_segment_need(segment_text, framework)
            selected_asset = None
            relevance = {"threshold": RELEVANCE_ACCEPT_THRESHOLD, "candidates": []}
            mode = "none"
            if score >= 0.34:
                mode, selected_asset, relevance = choose_mode(segment_text, slide_assets, framework)
            segment_decisions.append(
                {
                    "segment_index": idx,
                    "text": segment,
                    "need_score": round(score, 2),
                    "mode": mode,
                    "query": build_search_query(segment_text, framework) if mode == "search" else "",
                    "prompt": build_generation_prompt(segment, title, framework) if mode == "generate" else "",
                    "subject_framework": framework,
                    "selected_existing_asset": asdict(selected_asset) if selected_asset else None,
                    "existing_asset_relevance": relevance,
                }
            )

        actionable = next((item for item in segment_decisions if item["mode"] != "none"), None)
        rejection_summary = summarize_asset_rejections(segment_decisions, slide_assets)
        combined_text = f"{deck_context} {title} {' '.join(segments[:4])}"
        background_decor = {
            "mode": "search",
            "query": build_background_decor_query(combined_text, slide_framework),
            "purpose": "background_decor",
        } if should_plan_background_decor(slide_framework, combined_text) else {
            "mode": "none",
            "query": "",
            "purpose": "background_decor",
        }
        decision = {
            "slide_number": slide_num,
            "title": title,
            "recommended_mode": actionable["mode"] if actionable else "none",
            "reason": visual_decision_reason(actionable, slide_framework, rejection_summary),
            "subject_framework": slide_framework,
            "background_decor": background_decor,
            "relevance_threshold": RELEVANCE_ACCEPT_THRESHOLD,
            "existing_assets": [asdict(asset) for asset in slide_assets],
            "segments": segment_decisions,
        }
        decisions.append(decision)
    return decisions


def run_command(args: list[str], cwd: Path) -> bool:
    env = None
    if VENDOR_DIR.exists():
        env = os.environ.copy()
        existing = env.get("PYTHONPATH", "")
        env["PYTHONPATH"] = str(VENDOR_DIR) + ((";" + existing) if existing else "")
    try:
        subprocess.run(args, cwd=str(cwd), check=True, env=env)
        return True
    except subprocess.CalledProcessError:
        return False


def execute_decisions(
    project: Path,
    decisions: list[dict[str, Any]],
    *,
    allow_search: bool,
    allow_generate: bool,
    limit: int | None,
    background_only: bool = False,
    force_background_decor: bool = False,
) -> list[dict[str, Any]]:
    output_dir = project / "images" / "visual_assets"
    output_dir.mkdir(parents=True, exist_ok=True)
    decor_dir = project / "images" / "background_decor"
    decor_dir.mkdir(parents=True, exist_ok=True)
    completed = []
    count = 0

    for decision in decisions:
        if limit is not None and count >= limit:
            break
        slide_num = int(decision["slide_number"])

        background_plan = decision.get("background_decor") or {}
        if (
            allow_search
            and background_plan.get("mode") == "search"
            and (limit is None or count < limit)
        ):
            decor_filename = f"slide_{slide_num:02d}_background.jpg"
            decor_path = decor_dir / decor_filename
            decor_selected = {
                "slide_number": slide_num,
                "kind": "background_decor",
                "mode": "search",
                "status": "planned",
                "filename": decor_filename,
                "query": background_plan.get("query", ""),
            }
            if decor_path.exists() and decor_path.stat().st_size > 40_000 and not force_background_decor:
                decor_selected.update({"status": "exists", "path": str(decor_path)})
                completed.append(decor_selected)
                count += 1
            else:
                framework = decision.get("subject_framework") or {}
                query = background_plan.get("query") or build_background_decor_query(decision.get("title", ""), framework)
                cache_framework = dict(framework)
                cache_framework["cache_purpose"] = "background_decor"
                cache_framework["background_variant"] = f"slide_{slide_num:02d}"
                cache_hit = restore_visual_cache(query, cache_framework, decor_path, decor_dir)
                ok = bool(cache_hit)
                if not cache_hit:
                    ok = run_command(
                        [
                            sys.executable,
                            str(SCRIPT_DIR / "image_search.py"),
                            query,
                            "--filename",
                            decor_filename,
                            "--slide",
                            f"{slide_num:02d}",
                            "--purpose",
                            "background_decor",
                            "--orientation",
                            "landscape",
                            "-o",
                            str(decor_dir),
                        ],
                        cwd=SCRIPT_DIR.parent.parent,
                    )
                decor_selected.update({"status": "ready" if ok else "failed", "cache_hit": bool(cache_hit)})
                if ok and decor_path.exists():
                    source_item = _read_image_sources(decor_dir).get(decor_filename, {})
                    decor_selected["source"] = source_item
                    decor_selected["path"] = str(decor_path)
                    if cache_hit:
                        decor_selected["cache"] = cache_hit
                    else:
                        decor_selected["cache"] = store_visual_cache(
                            query,
                            cache_framework,
                            decor_path,
                            source_item,
                            {"accepted": True, "reason": "background decor asset"},
                        )
                completed.append(decor_selected)
                count += 1

        if limit is not None and count >= limit:
            break
        if background_only:
            continue
        mode = decision["recommended_mode"]
        if mode == "none":
            continue

        filename = f"slide_{slide_num:02d}_visual.jpg"
        selected = {
            "slide_number": slide_num,
            "kind": "visual_asset",
            "mode": mode,
            "status": "planned",
            "filename": filename,
        }

        if mode == "existing":
            segment = next((item for item in decision["segments"] if item["mode"] == mode), None)
            selected_asset = segment.get("selected_existing_asset") if segment else None
            if selected_asset and selected_asset.get("quality", {}).get("usable"):
                src = Path(selected_asset["path"])
                suffix = src.suffix.lower() or ".jpg"
                dst = output_dir / f"slide_{slide_num:02d}_existing{suffix}"
                shutil.copy2(src, dst)
                selected.update(
                    {
                        "status": "ready",
                        "filename": dst.name,
                        "path": str(dst),
                        "source_path": str(src),
                        "relevance": segment.get("existing_asset_relevance") if segment else {},
                    }
                )
                completed.append(selected)
                count += 1
            continue

        segment = next((item for item in decision["segments"] if item["mode"] == mode), None)
        if not segment:
            continue

        if mode == "search" and allow_search:
            query = segment["query"]
            framework = segment.get("subject_framework") or decision.get("subject_framework") or {}
            output_path = output_dir / filename
            if output_path.exists() and output_path.stat().st_size > 50_000:
                cache_hit = None
                ok = True
                selected.update({"status": "exists", "path": str(output_path), "query": query, "cache_hit": False})
            else:
                cache_hit = restore_visual_cache(query, framework, output_path, output_dir)
                ok = bool(cache_hit)
            if not output_path.exists() and not cache_hit:
                ok = run_command(
                    [
                        sys.executable,
                        str(SCRIPT_DIR / "image_search.py"),
                        query,
                        "--filename",
                        filename,
                        "--slide",
                        f"{slide_num:02d}",
                        "--purpose",
                        "visual_asset",
                        "--orientation",
                        "landscape",
                        "-o",
                        str(output_dir),
                    ],
                    cwd=SCRIPT_DIR.parent.parent,
                )
            selected.update(
                {
                    "status": "ready" if ok else "failed",
                    "query": query,
                    "cache_hit": bool(cache_hit),
                }
            )
            if ok:
                source_item = _read_image_sources(output_dir).get(filename, {})
                audit = audit_downloaded_visual(
                    output_path,
                    source_item,
                    query=query,
                    framework=framework,
                )
                selected["web_relevance_audit"] = audit.get("metadata", {})
                selected["image_content_audit"] = audit.get("content", {})
                selected["visual_audit"] = audit
                selected["source"] = source_item
                if cache_hit:
                    selected["cache"] = cache_hit
                if audit["accepted"] and output_path.exists():
                    selected["path"] = str(output_path)
                    if not cache_hit:
                        selected["cache"] = store_visual_cache(query, framework, output_path, source_item, audit)
                else:
                    selected["status"] = "rejected"
                    selected["rejected_path"] = _quarantine_rejected_asset(output_path) if output_path.exists() else ""
                    selected["reject_reason"] = audit["reason"]
            completed.append(selected)
            count += 1
            continue

        if mode == "generate" and allow_generate:
            stem = f"slide_{slide_num:02d}_generated"
            ok = run_command(
                [
                    sys.executable,
                    str(SCRIPT_DIR / "image_gen.py"),
                    segment["prompt"],
                    "--aspect_ratio",
                    "16:9",
                    "--image_size",
                    "1K",
                    "--filename",
                    stem,
                    "-o",
                    str(output_dir),
                ],
                cwd=SCRIPT_DIR.parent.parent,
            )
            generated = next(output_dir.glob(f"{stem}.*"), None)
            selected.update({"status": "ready" if ok and generated else "failed", "prompt": segment["prompt"]})
            if generated:
                content_audit = analyze_image_content(generated)
                selected["image_content_audit"] = content_audit
                if content_audit.get("accepted") is False:
                    selected["status"] = "rejected"
                    selected["rejected_path"] = _quarantine_rejected_asset(generated)
                    selected["reject_reason"] = content_audit.get("reason", "")
                else:
                    selected.update({"filename": generated.name, "path": str(generated)})
            completed.append(selected)
            count += 1

    return completed


def write_manifest(project: Path, payload: dict[str, Any]) -> Path:
    path = project / "images" / "visual_asset_manifest.json"
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("project_path", type=Path, help="PPT Master project directory")
    parser.add_argument("--execute", action="store_true", help="Run search/generation for planned assets")
    parser.add_argument("--no-search", action="store_true", help="Do not call image_search.py in execute mode")
    parser.add_argument("--no-generate", action="store_true", help="Do not call image_gen.py in execute mode")
    parser.add_argument("--limit", type=int, default=None, help="Limit executed assets")
    parser.add_argument("--background-only", action="store_true", help="Only search/download background decor assets")
    parser.add_argument("--force-background-decor", action="store_true", help="Overwrite existing background decor assets")
    parser.add_argument("--min-width", type=int, default=900, help="Minimum usable source image width")
    parser.add_argument("--min-height", type=int, default=500, help="Minimum usable source image height")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    project = args.project_path
    if not project.exists():
        raise FileNotFoundError(project)

    assets = extract_source_images(project, args.min_width, args.min_height)
    slides = load_slide_texts(project)
    decisions = plan_visual_assets(slides, assets, project.name)
    executed = []
    if args.execute:
        executed = execute_decisions(
            project,
            decisions,
            allow_search=not args.no_search,
            allow_generate=not args.no_generate,
            limit=args.limit,
            background_only=args.background_only,
            force_background_decor=args.force_background_decor,
        )

    payload = {
        "project": project.name,
        "generated_at": now_iso(),
        "mode": "execute" if args.execute else "dry_run",
        "quality_threshold": {
            "min_width": args.min_width,
            "min_height": args.min_height,
        },
        "web_relevance_threshold": WEB_RELEVANCE_ACCEPT_THRESHOLD,
        "visual_cache_dir": str(VISUAL_CACHE_DIR),
        "extracted_assets": [asdict(asset) for asset in assets],
        "decisions": decisions,
        "executed_assets": executed,
    }
    manifest_path = write_manifest(project, payload)

    print(f"Extracted source images: {len(assets)}")
    print(f"Planned slide decisions: {len(decisions)}")
    print(f"Executed assets: {len(executed)}")
    print(f"Manifest: {manifest_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
